#!/usr/bin/env python3
"""
pptx_xycut_demo.py
──────────────────
Extract editable text + native tables from a .pptx and segment them
with a recursive XY‑cut.  Output is JSON (reading order preserved).

$ pip install python-pptx
$ python pptx_xycut_demo.py slides.pptx > slides.json
"""
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import json
from typing import List, Dict, Any, Tuple

# ──────────────────── 1. recursive extractor ──────────────────────
def iter_text_and_tables(shapes, off=(0, 0)) -> List[Dict[str, Any]]:
    """Return a flat list of {'type','text'|'rows','bbox'} dicts."""
    out = []
    ox, oy = off
    for sh in shapes:
        x, y = ox + sh.left, oy + sh.top
        bbox = (x, y, x + sh.width, y + sh.height)

        # dive into groups
        if sh.shape_type == MSO_SHAPE_TYPE.GROUP:
            out.extend(iter_text_and_tables(sh.shapes, off=(x, y)))

        # ordinary text box / placeholder / auto‑shape
        elif sh.has_text_frame:
            txt = sh.text_frame.text.strip()
            if txt:
                out.append({"type": "text", "bbox": bbox, "text": txt})

        # native PowerPoint table
        elif sh.has_table:
            rows = [[cell.text.strip() for cell in row.cells]
                    for row in sh.table.rows]
            out.append({"type": "table", "bbox": bbox, "rows": rows})
    return out

# ──────────────────── 2. XY‑cut segmentation ──────────────────────
def xy_cut(blocks: List[Dict[str, Any]],
           direction: str = "h",
           thresh: float = 0.06,
           depth: int = 0,
           max_depth: int = 10) -> List[List[Dict[str, Any]]]:
    """
    Recursively split blocks on whitespace.  'direction' alternates:
    'h' = test horizontal gaps (i.e. split top/bottom),
    'v' = test vertical   gaps (split left/right).
    Returns a list of leaf‑region block lists in reading order.
    """
    if len(blocks) <= 1 or depth >= max_depth:
        return [blocks]

    # sort once in the test direction
    key = (1 if direction == "h" else 0)  # y0 or x0
    blocks = sorted(blocks, key=lambda b: b["bbox"][key])

    # measure gaps between adjacent blocks
    gaps: List[Tuple[float, int]] = []  # (gap_size, index_after)
    for i, (a, b) in enumerate(zip(blocks, blocks[1:]), start=1):
        if direction == "h":
            gap = b["bbox"][1] - a["bbox"][3]        # topB – bottomA
            ref = max(bk["bbox"][3] - bk["bbox"][1] for bk in blocks)
        else:
            gap = b["bbox"][0] - a["bbox"][2]        # leftB – rightA
            ref = max(bk["bbox"][2] - bk["bbox"][0] for bk in blocks)
        if gap / ref >= thresh:                      # significant gap?
            gaps.append((gap, i))

    if not gaps:
        return [blocks]                              # no split → leaf

    # pick the largest gap, split there
    _, split_idx = max(gaps, key=lambda t: t[0])
    left, right = blocks[:split_idx], blocks[split_idx:]

    next_dir = "v" if direction == "h" else "h"
    return  (xy_cut(left,  next_dir, thresh, depth+1, max_depth) +
             xy_cut(right, next_dir, thresh, depth+1, max_depth))

# ──────────────────── 3. end‑to‑end driver ────────────────────────
def process_pptx(path: str, gap_thresh: float = 0.06) -> List[Dict[str, Any]]:
    prs = Presentation(path)
    slides_json = []

    for idx, slide in enumerate(prs.slides, start=1):
        blocks = iter_text_and_tables(slide.shapes)

        # separate tables so XY‑cut only sequences *text* blocks
        text_blocks = [b for b in blocks if b["type"] == "text"]
        tables      = [b for b in blocks if b["type"] == "table"]

        # run XY‑cut (alternating h/v starting with horizontal)
        segments = xy_cut(text_blocks, "h", gap_thresh)

        # flatten leaf segments in produced reading order
        ordered = []
        for seg in segments:
            if not seg:
                continue
            if len(seg) == 1:
                ordered.append(seg[0])
            else:
                # merge contiguous shapes inside same segment
                merged_txt = " ".join(b["text"] for b in seg)
                bbox0 = min(b["bbox"][0] for b in seg)
                bboxt = min(b["bbox"][1] for b in seg)
                bbox1 = max(b["bbox"][2] for b in seg)
                bboxb = max(b["bbox"][3] for b in seg)
                ordered.append({"type": "text",
                                "bbox": (bbox0, bboxt, bbox1, bboxb),
                                "text": merged_txt})

        # insert tables into their geometric position in the order list
        for tbl in tables:
            inserted = False
            for i, blk in enumerate(ordered):
                if tbl["bbox"][1] < blk["bbox"][1]:  # table higher on slide
                    ordered.insert(i, tbl); inserted = True; break
            if not inserted:
                ordered.append(tbl)

        slides_json.append({"slide_index": idx, "blocks": ordered})
    return slides_json


if __name__ == "__main__":
    import sys, pathlib, pprint
    if len(sys.argv) < 2:
        sys.exit("Usage: python pptx_xycut_demo.py file.pptx")
    pptx = pathlib.Path(sys.argv[1])
    result = process_pptx(str(pptx))
    print(json.dumps(result, indent=2))
