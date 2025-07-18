from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml import parse_xml
from pptx.text.text import _Paragraph
from pathlib import Path
from collections import defaultdict
from statistics import median, quantiles
import statistics, math
import re, json, sys
import base64

EMU_PER_PT = 914400 / 72          # 1 pt = 1/72 in ; 1 in = 914 400 EMU
BIN_WIDTH  = 1                    # 1 pt vertical bin for gutters
COL_FACTOR = 1.0                  # ≥ 1 char‑width → candidate column gutter
ROW_FACTOR = 1.5                  # ≥ 1.5 line‑heights → candidate paragraph gap
BULLET_RE   = re.compile(r"^[\u2022\u2023\u25E6\u2043\u2219\u00B7]+$")
TEXT_SCALE  = 0.52


def _font_size(par: _Paragraph, default=12.0) -> float:
    for run in par.runs:
        if run.font.size:
            return run.font.size.pt
    if par.font.size:
        return par.font.size.pt
    return default

def _bullet_char(par: _Paragraph) -> str | None:
    ppr = parse_xml(par._pPr.xml) if par._pPr is not None else None
    if ppr is None:
        return None

    if ppr.find('.//a:buNone', ppr.nsmap) is not None:
        return None

    bu_char = ppr.find('.//a:buChar', ppr.nsmap)
    if bu_char is not None and bu_char.get("char"):
        return bu_char.get("char")

    return "•"

def _indent_pts(par):
    fmt = getattr(par, "paragraph_format", None)
    if fmt is not None and fmt.left_indent is not None:
        try:
            return fmt.left_indent.pt
        except AttributeError:
            return fmt.left_indent / EMU_PER_PT

    ppr = parse_xml(par._pPr.xml) if par._pPr is not None else None
    if ppr is not None and ppr.get("marL") is not None:
        return int(ppr.get("marL")) / EMU_PER_PT
    return 0.0

def _shape_table_block(sh):
    tbl = sh.table

    n_rows = getattr(tbl, "row_count", len(tbl.rows))
    n_cols = getattr(tbl, "column_count", len(tbl.columns))

    rows = [
        [
            tbl.cell(r, c).text_frame.text.strip()
            for c in range(n_cols)
        ]
        for r in range(n_rows)
    ]

    x0 = sh.left  / EMU_PER_PT
    y0 = sh.top   / EMU_PER_PT
    x1 = x0 + sh.width  / EMU_PER_PT
    y1 = y0 + sh.height / EMU_PER_PT
    return {
        "type": "table",
        "bbox": [x0, y0, x1, y1],
        "rows": rows,
    }

def _column_props(text_frame):
    fmt = getattr(text_frame, "text_frame_format", None)
    if fmt is not None:
        cnt = fmt.column_count or 1
        spc = (fmt.column_spacing.pt if fmt.column_spacing else 0)
        return cnt, spc

    bodyPr = text_frame._element.bodyPr
    cnt = int(bodyPr.get("numCol", "1"))
    spc = int(bodyPr.get("spcCol", "0")) / EMU_PER_PT
    return max(1, cnt), spc

def _column_width(shape, text_frame):
    cnt, spc = _column_props(text_frame)
    return (shape.width / EMU_PER_PT - (cnt - 1) * spc) / cnt, cnt, spc

def extract_rows_from_slide(slide):
    rows, table_blocks, image_blocks = [], [], []

    body_sizes = []
    for sh in slide.shapes:
        if sh.has_text_frame:
            body_sizes.extend(_font_size(p) for p in sh.text_frame.paragraphs)
    body_med = statistics.median(body_sizes) if body_sizes else 12.0

    for sh in slide.shapes:
        if sh.has_table:
            table_blocks.append(_shape_table_block(sh))
            continue
        if sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
            pic_blk = picture_block(sh, slide.part.slide_width / EMU_PER_PT,
                                        slide.part.slide_height / EMU_PER_PT)
            if pic_blk:
                image_blocks.append(pic_blk)
            continue           
        if not sh.has_text_frame:
            continue
        tf = sh.text_frame

        col_w, col_cnt, gutter = _column_width(sh, tf)
        base_x = sh.left / EMU_PER_PT
        base_y = sh.top  / EMU_PER_PT

        paras      = list(tf.paragraphs)
        total_para = len(paras)
        paras_per_col = math.ceil(total_para / col_cnt)

        y_cursor = base_y
        for para_idx, par in enumerate(paras):
            raw = par.text.replace("\xa0", " ").strip()
            if not raw and not _bullet_char(par):
                y_cursor += _font_size(par) * 1.15
                continue

            size   = _font_size(par)
  
            indent  = _indent_pts(par)
            bullet  = _bullet_char(par)

            if bullet and indent < 2 and size <= body_med:
                bullet = None

            text = (f"{bullet} {raw}"
                    if bullet and (not raw or raw[0] != bullet)
                    else raw or bullet)

            col_idx  = min(para_idx // paras_per_col, col_cnt - 1)
            x0       = base_x + col_idx * (col_w + gutter) + indent
            est_w    = max(1, len(text)) * size * TEXT_SCALE
            x1       = min(x0 + est_w, base_x + (col_idx + 1) * col_w)

            line_h = size * 1.15
            y0, y1 = y_cursor, y_cursor + line_h
            y_cursor += line_h

            rows.append({
                "x0": x0, "y0": y0, "x1": x1, "y1": y1,
                "size": size,
                "text": text,
                "baseline": y0,
                "is_bullet_only": BULLET_RE.fullmatch(text.strip()) is not None
            })

    rows.sort(key=lambda r: (r["y0"], r["x0"]))

    return rows, table_blocks, image_blocks

def char_widths(row_list):
    w = []
    for r in row_list:
        n = sum(1 for c in r["text"] if not c.isspace())
        if n:
            w.append((r["x1"] - r["x0"]) / n)
    return w

def compute_dominant_sizes(rows):
    cw, lh = [], []
    for r in rows:
        n = sum(1 for c in r["text"] if not c.isspace())
        if n:
            cw.append((r["x1"] - r["x0"]) / n)
        lh.append(r["y1"] - r["y0"])
    char = median(cw) if cw else 0
    if not lh:
        return char, 0
    q1 = quantiles(lh, n=4)[0]
    core = [h for h in lh if h >= q1]
    return char, median(core)

def gap(a, b):
    return max(0.0, b["y0"] - a["y1"])

def looks_like_heading(row, med_font, size_ratio=1.15, short_len=60):
    txt = row["text"].strip()
    if not any(ch.isalpha() for ch in txt):
        return False
    return row["size"] >= med_font * size_ratio and len(txt) <= short_len

def is_large_gap(g, line_h, med_gap):
    return g >= max(ROW_FACTOR * line_h, 2.5 * med_gap if med_gap else 0)

def find_vertical_gutter(rows, idx, page_w, char_w, tol=0.02):
    rx0 = min(rows[i]["x0"] for i in idx)
    rx1 = max(rows[i]["x1"] for i in idx)
    w   = rx1 - rx0
    if char_w == 0 or w <= 0:
        return None
    bins = int(w // BIN_WIDTH) + 1
    hist = [0]*bins
    for i in idx:
        x0 = max(rows[i]["x0"], rx0); x1 = min(rows[i]["x1"], rx1)
        b0 = int((x0-rx0)//BIN_WIDTH); b1 = int((x1-rx0)//BIN_WIDTH)
        for b in range(b0, b1+1):
            hist[b] += 1
    max_occ = max(1, int(tol*len(idx)))
    run_req = max(1, int((COL_FACTOR*char_w)//BIN_WIDTH))
    best = None; cur = None
    for j, occ in enumerate(hist+[max_occ+1]):
        if occ <= max_occ:
            cur = j if cur is None else cur
        elif cur is not None:
            run = j - cur
            if run >= run_req and (best is None or run > best[2]):
                best = (cur, j-1, run)
            cur = None
    if not best:
        return None
    g0, g1, _ = best
    return rx0 + g0*BIN_WIDTH, rx0 + g1*BIN_WIDTH + BIN_WIDTH

def find_horizontal_gap(rows, idx, page_h, line_h, *, tol=0.02):
    ry0 = min(rows[i]["y0"] for i in idx)
    ry1 = max(rows[i]["y1"] for i in idx)
    h   = ry1 - ry0
    if h <= 0 or line_h == 0:
        return None
    BIN_HEIGHT = 1
    bins = int(h // BIN_HEIGHT) + 1
    hist = [0]*bins
    for i in idx:
        y0 = max(rows[i]["y0"], ry0); y1 = min(rows[i]["y1"], ry1)
        b0 = int((y0-ry0)//BIN_HEIGHT); b1 = int((y1-ry0)//BIN_HEIGHT)
        for b in range(b0, b1+1):
            hist[b] += 1
    max_occ = max(1, int(tol*len(idx)))
    run_req = max(1, int((ROW_FACTOR*line_h)//BIN_HEIGHT))
    best = None; cur = None
    for j, occ in enumerate(hist+[max_occ+1]):
        if occ <= max_occ:
            cur = j if cur is None else cur
        elif cur is not None:
            run = j - cur
            if run >= run_req and (best is None or run > best[2]):
                best = (cur, j-1, run)
            cur = None
    if not best:
        return None
    g0, g1, _ = best
    return ry0 + g0*BIN_HEIGHT, ry0 + g1*BIN_HEIGHT + BIN_HEIGHT

def picture_block(pic, slide_w_pt, slide_h_pt):
    x0 = pic.left  / EMU_PER_PT
    y0 = pic.top   / EMU_PER_PT
    w  = pic.width  / EMU_PER_PT
    h  = pic.height / EMU_PER_PT
    x1, y1 = x0 + w, y0 + h

    if w < 40 and h < 40:
        return None
    if (w / max(1, h) > 10) or (h / max(1, w) > 10):
        return None
    if w >= 0.9 * slide_w_pt and h >= 0.9 * slide_h_pt:
        return None

    b64 = base64.b64encode(pic.image.blob).decode()
    return {
        "type":   "image",
        "bbox":   [round(x0, 2), round(y0, 2), round(x1, 2), round(y1, 2)],
        "width":  int(w),
        "height": int(h),
        "text":   b64,
    }


# ───────────────────────────────
#  Recursive XY‑cut
# ───────────────────────────────
def xy_cut_region(idx, rows, page_w, page_h,
                  *, min_block_w=0.02, min_block_h=0.02):
    if not idx or len(idx) == 1:
        return [idx]

    char_w, line_h = compute_dominant_sizes([rows[i] for i in idx])
    if char_w == 0 or line_h == 0:
        return [idx]

    # 1) vertical
    gutter = find_vertical_gutter(rows, idx, page_w, char_w)
    if gutter:
        gx0, gx1 = gutter
        left  = [i for i in idx if rows[i]["x1"] <= gx0]
        right = [i for i in idx if rows[i]["x0"] >= gx1]
        bridge = [i for i in idx if i not in left and i not in right
                  and rows[i]["x0"] < gx1 and rows[i]["x1"] > gx0]
        if left and right:
            min_w = min_block_w * page_w
            lw = max(rows[i]["x1"] for i in left) - min(rows[i]["x0"] for i in left)
            rw = max(rows[i]["x1"] for i in right) - min(rows[i]["x0"] for i in right)
            if lw >= min_w and rw >= min_w:
                return (xy_cut_region(sorted(left, key=lambda i: rows[i]["y0"]),
                                      rows, page_w, page_h)
                     + xy_cut_region(sorted(bridge, key=lambda i: rows[i]["y0"]),
                                      rows, page_w, page_h)
                     + xy_cut_region(sorted(right, key=lambda i: rows[i]["y0"]),
                                      rows, page_w, page_h))

    # 2) horizontal
    hgap = find_horizontal_gap(rows, idx, page_h, line_h)
    if hgap:
        gy0, gy1 = hgap
        upper = [i for i in idx if rows[i]["y1"] <= gy0]
        lower = [i for i in idx if rows[i]["y0"] >= gy1]
        if upper and lower:
            min_h = min_block_h * page_h
            uh = max(rows[i]["y1"] for i in upper) - min(rows[i]["y0"] for i in upper)
            lh = max(rows[i]["y1"] for i in lower) - min(rows[i]["y0"] for i in lower)
            if uh >= min_h and lh >= min_h:
                return (xy_cut_region(sorted(upper, key=lambda i: rows[i]["y0"]),
                                      rows, page_w, page_h)
                     + xy_cut_region(sorted(lower, key=lambda i: rows[i]["y0"]),
                                      rows, page_w, page_h))

    # 3) heading‑boosted gap analyzer
    by_top = sorted(idx, key=lambda i: rows[i]["y0"])
    med_font = statistics.median(rows[i]["size"] for i in idx)
    gaps_white, gaps_base = [], []
    for a, b in zip(by_top, by_top[1:]):
        white = gap(rows[a], rows[b])
        base  = rows[b]["y0"] - rows[a]["y0"]
        if looks_like_heading(rows[a], med_font) or looks_like_heading(rows[b], med_font):
            base = max(base, ROW_FACTOR * line_h + 1)
        gaps_white.append(white)
        gaps_base.append(base)
    med_gap = statistics.median(gaps_white) if gaps_white else 0
    big = [k for k, g in enumerate(gaps_base)
           if is_large_gap(g, line_h, med_gap)]
    if big:
        cut = max(big, key=lambda k: gaps_base[k]) + 1
        upper, lower = by_top[:cut], by_top[cut:]
        return (xy_cut_region(upper, rows, page_w, page_h)
             + xy_cut_region(lower, rows, page_w, page_h))

    return [idx]

def merge_bullets(rows, char_tol_factor=1.0):
    merged = []
    by_base = defaultdict(list)
    for r in rows:
        by_base[r["baseline"]].append(r)
    for base, seq in by_base.items():
        seq.sort(key=lambda r: r["x0"])
        pool = [
            (r["x1"]-r["x0"])/max(1,len(r["text"]))
            for r in seq if not r["is_bullet_only"]
        ]
        cw = median(pool) if pool else 1.0
        keep, skip = [], set()
        for i,r in enumerate(seq):
            if r["is_bullet_only"] and i+1 < len(seq):
                nxt = seq[i+1]
                if (abs(nxt["baseline"] - r["baseline"]) <= 3.0
                    and abs(nxt["x0"] - r["x0"]) <= char_tol_factor*cw):
                    nxt["text"] = f"{r['text']} {nxt['text']}"
                    nxt["x0"]   = min(r["x0"], nxt["x0"])
                    skip.add(id(r))
        for r in seq:
            if id(r) not in skip:
                keep.append(r)
        merged.extend(keep)
    merged.sort(key=lambda r:(r["y0"],r["x0"]))
    return merged

def make_output_chunk(seg, rows):
    x0 = min(rows[i]["x0"] for i in seg)
    y0 = min(rows[i]["y0"] for i in seg)
    x1 = max(rows[i]["x1"] for i in seg)
    y1 = max(rows[i]["y1"] for i in seg)
    font = statistics.median(rows[i]["size"] for i in seg)
    txt  = " ".join(rows[i]["text"] for i in seg)
    return {"type":"text","font_size":round(font,2),
            "bbox":[x0,y0,x1,y1],"text":txt}

def iterate_chunks_slide(slide, slide_w_pt, slide_h_pt):
    rows, tbl_blocks, img_blocks = extract_rows_from_slide(slide)
    rows             = merge_bullets(rows)

    idx   = list(range(len(rows)))
    segs  = xy_cut_region(idx, rows, slide_w_pt, slide_h_pt)
    text_blocks = [make_output_chunk(seg, rows) for seg in segs if seg]

    return img_blocks + tbl_blocks + text_blocks

def main(pptx_path: str):
    prs = Presentation(pptx_path)
    slide_w = prs.slide_width  / EMU_PER_PT
    slide_h = prs.slide_height / EMU_PER_PT
    result = []
    for n, slide in enumerate(prs.slides,1):
        blocks = iterate_chunks_slide(slide, slide_w, slide_h)
        result.append({"page": n, "blocks": blocks})
    print(json.dumps(result, indent=2, ensure_ascii=False))

def extract_blocks(pptx_path: str | Path):
    prs = Presentation(pptx_path)
    slide_w = prs.slide_width  / EMU_PER_PT
    slide_h = prs.slide_height / EMU_PER_PT

    for page_no, slide in enumerate(prs.slides, 1):
        for blk in iterate_chunks_slide(slide, slide_w, slide_h):
            if "page" not in blk:
                blk = {**blk, "page": page_no}
            yield blk

if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("usage: pptxycut.py <file.pptx>")
        sys.exit(1)
    main(sys.argv[1])